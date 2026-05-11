import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import io

def convert_pdf_to_pptx(pdf_path, pptx_path):
    try:
        # Create presentation object
        prs = Presentation()
        
        # Remove the default first slide if necessary (usually prs starts empty)
        # but standard python-pptx starts with 0 slides.
        
        # Open the PDF
        doc = fitz.open(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Convert page to image (pixmap)
            # Use a matrix to increase resolution (e.g., 2.0x zoom)
            zoom = 1.5
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Save pixmap to a memory buffer as PNG
            img_data = pix.tobytes("png")
            img_stream = io.BytesIO(img_data)
            
            # Add a blank slide (layout 6 is usually blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Set slide size to match page aspect ratio
            # PowerPoint default is 10x5.625 (16:9)
            # We'll stretch the image to fill the slide, or we could adjust prs.slide_width/height
            
            # For simplicity, we fill the slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add image to slide at (0,0) covering the whole slide
            slide.shapes.add_picture(img_stream, 0, 0, width=slide_width, height=slide_height)
            
        # Save the presentation
        prs.save(pptx_path)
        doc.close()
        print("SUCCESS")
    except Exception as e:
        print(f"ERROR: {e}")
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdfToPpt.py <input_pdf> <output_pptx>")
        sys.exit(1)

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    
    convert_pdf_to_pptx(input_pdf, output_pptx)
